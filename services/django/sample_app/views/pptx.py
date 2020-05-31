import json
import time
import datetime

from django.http import HttpResponse
from django.shortcuts import render
from django.template import RequestContext, loader

from pptx import Presentation
from pptx.util import Inches

from ..models.campaign_framework import *


def index(request):
    """

    :param request:
    :type request:
    :return:
    :rtype:
    """
    cf_step = CampaignFrameworkStep.objects.get(
        pk=5
    )

    filename = "OmniTemplateTest.pptx"
    prs = Presentation(filename)

    slide = prs.slides.get(prs.slides[0].slide_id)
    slide.shapes.title.text = "HI THERE"
    slide.placeholders[1].text = " fgdfg fd"

    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.text = "f esfsf dsfsdfds"

    """
    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('%d, %s' % (phf.idx, phf.type))
    """

    #for x in prs.slides:
        #print (x.slide_id)

    """
    ppt_layout = prs.slide_layouts[0]
    ppt_layout1 = prs.slide_layouts[1]
    ppt_layout2 = prs.slide_layouts[2]
    ppt_layout3 = prs.slide_layouts[3]
    ppt_layout4 = prs.slide_layouts[4]
    ppt_layout5 = prs.slide_layouts[5]
    ppt_layout6 = prs.slide_layouts[6]
    ppt_layout7 = prs.slide_layouts[7]
    ppt_layout8 = prs.slide_layouts[8]
    ppt_layout9 = prs.slide_layouts[9]
    ppt_layout10 = prs.slide_layouts[10]

    prs.slides.add_slide(ppt_layout)
    prs.slides.add_slide(ppt_layout1)
    prs.slides.add_slide(ppt_layout2)
    prs.slides.add_slide(ppt_layout3)
    prs.slides.add_slide(ppt_layout4)
    prs.slides.add_slide(ppt_layout5)
    prs.slides.add_slide(ppt_layout6)
    prs.slides.add_slide(ppt_layout7)
    prs.slides.add_slide(ppt_layout8)
    prs.slides.add_slide(ppt_layout9)
    prs.slides.add_slide(ppt_layout10)
    """

    #blank_layout = prs.slide_layouts[6]
    #slide = prs.slides.add_slide(blank_layout)
    top = 1.5
    width = 2
    height = 1
    for row in cf_step.json_layout.get('rows'):
        left = 0.5
        for col in row.get('columns'):
            label = col.get('label')
            value = col.get('name')
            text_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height)
            )
            text_box.text_frame.text = label + "\r" + str(value)
            left = left + 3

        top = top + 0.5

    prs.save(filename)

    data = {}
    template = loader.get_template("sample_app/save_pptx.html")
    return HttpResponse(template.render(data, request))


def create_pptx(request):
    """

    :param request:
    :type request:
    :return:
    :rtype:
    """
    prs = Presentation('test.pptx')
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    for x in slide.placeholders:
        print('%d %s' % (x.placeholder_format.idx, x.name))

    print("=======================")

    for shape in slide.shapes:
        print('%s' % shape.shape_type)

    print ("=======================")

    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('%d, %s' % (phf.idx, phf.type))

    #return HttpResponse('hi')

    #title_slide_layout = prs.slide_layouts[0]
    #slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')

    data = {}
    return HttpResponse(
        json.dumps(data),
        content_type="application/json"
    )
