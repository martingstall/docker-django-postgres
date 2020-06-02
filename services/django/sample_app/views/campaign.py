import json
import requests
import time
import datetime

from django.core.files.storage import FileSystemStorage
from django.http import HttpResponse
from django.shortcuts import render
from django.template import RequestContext, loader
from django.views.decorators.csrf import csrf_exempt

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches

from ..models.campaign import *
from ..models.campaign_framework import *

from ..data_connections.factory import Factory as DataConnFactory


def campaign_details(request, campaign_id):
    """

    :param request:
    :type request:
    :param campaign_id:
    :type campaign_id:
    :return:
    :rtype:
    """
    campaign = Campaign.objects.get(pk=campaign_id)
    steps = CampaignFrameworkStep.objects.filter(
        phase_id__in=CampaignFrameworkPhase.objects.filter(
            campaign_framework_id = campaign.campaign_framework_id
        )
    )

    data = {
        "campaign": campaign,
        "steps": steps
    }
    template = loader.get_template("sample_app/campaign_detail.html")
    return HttpResponse(template.render(data, request))


def view_step(request, campaign_id, cf_step_id):
    """

    :param request:
    :type request:
    :param campaign_id:
    :type campaign_id:
    :param cf_step_id:
    :type cf_step_id:
    :return:
    :rtype:
    """
    campaign = Campaign.objects.get(pk=campaign_id)

    cf_step = CampaignFrameworkStep.objects.get(
        pk=cf_step_id
    )

    campaign_step = CampaignStepData.objects.get(
        campaign_id=campaign.id,
        campaign_framework_step_id=cf_step.id
    )
    campaign_step_data = campaign_step.campaign_step_data

    previous_steps = CampaignFrameworkStep.objects.values_list('pk', flat=True).filter(
        phase_id__in=CampaignFrameworkPhase.objects.only('pk').filter(
            campaign_framework_id=campaign.campaign_framework_id
        ),
        order__lt=cf_step.order,
        json_layout__contains={"rows":[{"columns":[{"api":{}}]}]}
    )
    persisted_objects = CampaignStepData.objects.values('campaign_step_data').filter(
        pk__in=previous_steps
    )

    data = {
        "campaign": campaign,
        "cf_step": cf_step,
        "json_layout": json.dumps(cf_step.json_layout),
        "campaign_step_data": json.dumps(campaign_step_data),
        "persisted_objects": json.dumps(list(persisted_objects)),
        "which_normalization_example": campaign_step_data.get('app_1')
    }
    template = loader.get_template("sample_app/" + cf_step.html_layout)
    return HttpResponse(template.render(data, request))


@csrf_exempt
def save_step_data(request, campaign_id, cf_step_id):
    """

    :param request:
    :type request:
    :param campaign_id:
    :type campaign_id:
    :param cf_step_id:
    :type cf_step_id:
    :return:
    :rtype:
    """
    data = {
        "status": "success"
    }

    try:
        form_data = json.dumps(request.POST)
        step = CampaignStepData.objects.get(
            campaign_id=campaign_id,
            campaign_framework_step_id=cf_step_id
        )
        step.campaign_step_data = json.loads(form_data)
        step.save()

    except Exception as e:
        data = {
            "status": "fail",
            "reason": str(e)
        }

    return HttpResponse(
        json.dumps(data),
        content_type="application/json"
    )


def create_campaign_pptx(request, campaign_id):
    """

    :param request:
    :type request:
    :param campaign_id:
    :type campaign_id:
    :return:
    :rtype:
    """
    campaign = Campaign.objects.get(pk=campaign_id)
    cf_steps = CampaignFrameworkStep.objects.filter(
        phase_id__in=CampaignFrameworkPhase.objects.filter(
            campaign_framework_id=campaign.campaign_framework_id
        )
    )

    filename = "OmniTemplateTest.pptx"
    prs = Presentation(filename)

    intro_slide = prs.slides.get(prs.slides[0].slide_id)
    intro_slide.shapes.title.text = campaign.name
    intro_slide.placeholders[1].text = "Campaign description here..."

    slide_template = prs.slide_layouts[8]
    for cf_step in cf_steps:
        slide = prs.slides.add_slide(slide_template)
        slide.shapes.title.text = cf_step.name

        campaign_step = CampaignStepData.objects.get(
            campaign_id=campaign.id,
            campaign_framework_step_id=cf_step.id
        )
        campaign_step_data = campaign_step.campaign_step_data

        top = 1.5
        height = 1
        for row in cf_step.json_layout.get('rows'):
            left = 0.5
            for col in row.get('columns'):
                try:
                    width = int(col.get('class').replace('col-sm-', ''))
                except Exception as e:
                    print ("Couldn't determine width: ", e)
                    width = 1
                label = col.get('label')
                value = campaign_step_data[col.get('name')]
                text_box = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                text_box.text_frame.text = label + ": " + str(value)
                left = left + width

            top = top + 0.5

    """
    slide = prs.slides.add_slide(slide_template)
    chart_data = ChartData()
    chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
    chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    """

    prs.save(campaign.name + ".pptx")

    return HttpResponse("Created")


def step_api_call(request, campaign_id, cf_step_id):
    """

    :param request:
    :type request:
    :param campaign_id:
    :type campaign_id:
    :param cf_step_id:
    :type cf_step_id:
    :return:
    :rtype:
    """
    params = request.GET
    app = {
        "value": params.get('return[value]'),
        "label": params.get('return[label]')
    }

    from ..views.api_calls import get_audience_explorer_conn
    base_url = "https://audience.annalect.com/api"
    url = base_url + "/audience/list/01c22232-76d6-11e9-8de8-0a10e129e67c/449203c6-7ed3-11e8-8b6b-0a35455287ac"
    response = requests.get(
        url,
        cookies=get_audience_explorer_conn(request)
    )
    app_data = json.loads(response.text)

    DataClass = DataConnFactory.getDataClass(params.get('system_key'))
    dataClass = DataClass()
    data = dataClass.normalize_select_element(app, app_data)

    return HttpResponse(
        json.dumps(data),
        content_type="application/json"
    )


def upload_step_html_file(request, campaign_id, cf_step_id):
    """

    :param request:
    :type request:
    :param campaign_id:
    :type campaign_id:
    :param step_id:
    :type step_id:
    :return:
    :rtype:
    """
    if request.method == "POST" and request.FILES["html_file"]:
        html_file = request.FILES["html_file"]
        fs = FileSystemStorage()
        filename = fs.save(html_file.name, html_file)
        html_file_url = fs.url(filename)

    data = {

    }
    template = loader.get_template("sample_app/upload_step_html_file.html")
    return HttpResponse(template.render(data, request))


def parse_html(request):
    """

    :param request:
    :type request:
    :return:
    :rtype:
    """
    fs = FileSystemStorage()
    file = fs.open('uploaded_step_example.html')
    data = file.read().decode('utf-8')
    data = data.strip()
    data = data.replace('\n', '')
    data = data.replace('\r', '')

    json = {
        "rows": []
    }
    html_rows = data.split('<div class="row">')
    for html_row in html_rows:
        json_row = {
            "columns": []
        }
        # print(html_row)
        # print("")
        # print("===================")
        html_cols = html_row.split('<div class="col')
        for html_col in html_cols:
            if '-sm-' not in html_col:
                continue

            json_col = {
                'class': 'col-' + str(html_col)
            }
            print(html_col)
            print("")
            print("===================")
            json_row['columns'].append(json_col)

        json['rows'].append(json_row)

    print("")
    print("")
    print(json)
    print("")
    print("")


def lit_table_example(request):
    """

    :param request:
    :type request:
    :return:
    :rtype:
    """
    data = {

    }
    template = loader.get_template("sample_app/lit_table_example.html")
    return HttpResponse(template.render(data, request))
